#!/usr/bin/env python3
"""
PPTX Packager - Combine generated slide images into a .pptx file.

Takes a directory of slide images and packages them into a PowerPoint
presentation with 16:9 aspect ratio, full-bleed images.

Usage:
  python skills/nano-banana/scripts/package_pptx.py \
    --dir ppt_output \
    --output my_presentation.pptx
"""

import argparse
import os
import sys
from pathlib import Path


def package_pptx(image_dir: str, output_path: str) -> str:
    """
    Package slide images into a .pptx file.

    Args:
        image_dir: Directory containing slide-XX.png images.
        output_path: Output .pptx file path.

    Returns:
        Path to the saved .pptx file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("Error: python-pptx not installed")
        print("Please run: pip install python-pptx")
        sys.exit(1)

    # Find all slide images, sorted by number
    images = sorted(
        [f for f in os.listdir(image_dir) if f.startswith("slide-") and f.endswith(".png")]
    )

    if not images:
        print(f"FAIL: No slide images found in {image_dir}/")
        sys.exit(1)

    # Create presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for img_name in images:
        img_path = os.path.join(image_dir, img_name)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height,
        )

    # Ensure output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    prs.save(output_path)
    print(f"OK: {output_path} ({len(images)} slides)")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Package slide images into a .pptx file",
        epilog="""
Example usage:
  python skills/nano-banana/scripts/package_pptx.py \\
    --dir ppt_output/images \\
    --output presentation.pptx
""",
    )
    parser.add_argument("--dir", required=True, help="Directory containing slide-XX.png images")
    parser.add_argument("--output", required=True, help="Output .pptx file path")

    args = parser.parse_args()

    if not os.path.isdir(args.dir):
        print(f"Error: {args.dir} not found")
        sys.exit(1)

    package_pptx(args.dir, args.output)


if __name__ == "__main__":
    main()
